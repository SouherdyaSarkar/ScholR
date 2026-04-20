import os
import sys
import io
import fitz
import docx
import pptx
from fastapi import UploadFile
config_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../utils/Python config"))
sys.path.append(config_dir)
from supabaseAdminConfig import supabase

class FileHandler:

    async def extract_text(self, file: UploadFile) -> str:
        extension = file.filename.split('.')[-1].lower()
        content = await file.read()
        extracted_text = ""

        try:
            if extension == 'pdf':
                extracted_text = self._extract_from_pdf(content)
            elif extension == 'docx':
                extracted_text = self._extract_from_docx(content)
            elif extension == 'pptx':
                extracted_text = self._extract_from_pptx(content)
            elif extension in ['doc', 'ppt']:
                extracted_text = f"Error: Legacy format .{extension} is not supported. Please use .docx or .pptx."
            else:
                extracted_text = "Error: Unsupported file format."
        except Exception as e:
            print(f"Error extracting text from {file.filename}: {e}")
            extracted_text = ""
            
        return extracted_text

    def _extract_from_pdf(self, content: bytes) -> str:
        text = ""
        with fitz.open(stream=content, filetype="pdf") as doc:
            for page in doc:
                text += page.get_text() + "\n"
        return text

    def _extract_from_docx(self, content: bytes) -> str:
        text = ""
        doc = docx.Document(io.BytesIO(content))
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text

    def _extract_from_pptx(self, content: bytes) -> str:
        text = ""
        presentation = pptx.Presentation(io.BytesIO(content))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    async def uploadFile(self, userID: str, topic : str, file: UploadFile):
        try:
            extracted_text = await self.extract_text(file)
            print(f"Extracted {len(extracted_text)} characters from {file.filename}")
            content_type_map = {
            ".pdf": "application/pdf",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }

            ext = os.path.splitext(file.filename)[1].lower()
            content_type = content_type_map.get(ext, "application/octet-stream")
            await file.seek(0)
            file_bytes = await file.read()
            res = supabase.storage.from_('scholr').upload(
                file.filename, 
                file_bytes,
                file_options={"content-type": content_type}
            )
            res_json = res.data
            print("Upload Response:", res_json)
            
            db_res = supabase.table("documents").insert({
            "userID": userID,
            "topic": topic,
            "filename": file.filename,
            "filepath": res_json.path,
        }).execute()
            
            return {
                "filename": file.filename, 
                "text_length": len(extracted_text), 
                "status": "success",
                "extracted_text": extracted_text[:100] + "..." # return a snippet for debugging
            }

        except Exception as e:
            print(f"Error in uploadFile: {e}")
            return {"status": "error", "message": str(e)}

    
file_handler = FileHandler()